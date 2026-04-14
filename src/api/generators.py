from fpdf import FPDF
import os
from datetime import datetime

class PDFGenerator:
    """
    Generates professional PDF Knowledge Articles from extracted video data.
    """
    def __init__(self, output_dir="data/outputs"):
        self.output_dir = output_dir
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

    def generate(self, video_title, youtube_url, richness_score, concepts):
        """
        Creates a PDF file.
        :param concepts: List of dicts with 'text', 'synthesized_article', 'timestamp'.
        """
        pdf = FPDF()
        pdf.add_page()
        
        # Title
        pdf.set_font("Arial", 'B', 24)
        pdf.cell(0, 20, "Knowledge Article (V2K Sentinel)", 0, 1, 'C')
        
        # Header Info
        pdf.set_font("Arial", '', 12)
        pdf.cell(0, 10, f"Source: {video_title}", 0, 1)
        pdf.cell(0, 10, f"URL: {youtube_url}", 0, 1)
        pdf.cell(0, 10, f"Richness Score: {richness_score:.2f}/1.00", 0, 1)
        pdf.cell(0, 10, f"Generated on: {datetime.now().strftime('%Y-%m-%d %H:%M')}", 0, 1)
        pdf.ln(10)
        
        # TOC
        pdf.set_font("Arial", 'B', 16)
        pdf.cell(0, 10, "Table of Contents", 0, 1)
        pdf.set_font("Arial", '', 12)
        for i, concept in enumerate(concepts):
            pdf.cell(0, 8, f"{i+1}. Segment at {concept['timestamp']:.2f}s", 0, 1)
        pdf.add_page()

        # Detailed Content
        for i, concept in enumerate(concepts):
            pdf.set_font("Arial", 'B', 14)
            pdf.cell(0, 10, f"Segment {i+1} (at {concept['timestamp']:.2f}s)", 0, 1)
            pdf.ln(2)
            
            # Transcript Summary
            pdf.set_font("Arial", 'I', 10)
            pdf.multi_cell(0, 6, f"Transcript: {concept['text'][:300]}...")
            pdf.ln(5)
            
            # Synthesized Knowledge
            pdf.set_font("Arial", '', 11)
            pdf.multi_cell(0, 7, concept.get('synthesized_article', 'No synthesis available.'))
            pdf.ln(10)

        safe_title = "".join([c for c in video_title if c.isalnum() or c in (' ', '_')]).rstrip()
        filename = f"{safe_title}_Knowledge_Article.pdf"
        filepath = os.path.join(self.output_dir, filename)
        pdf.output(filepath)
        return filepath

class MarkdownGenerator:
    """
    Generates GFM (GitHub Flavored Markdown) summaries.
    """
    def __init__(self, output_dir="data/outputs"):
        self.output_dir = output_dir
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

    def generate(self, video_title, youtube_url, richness_score, concepts):
        safe_title = "".join([c for c in video_title if c.isalnum() or c in (' ', '_')]).rstrip()
        filename = f"{safe_title}.md"
        filepath = os.path.join(self.output_dir, filename)
        
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(f"# {video_title}\n\n")
            f.write(f"- **URL:** {youtube_url}\n")
            f.write(f"- **Richness Score:** {richness_score}\n\n")
            f.write("## Table of Contents\n")
            for i, c in enumerate(concepts):
                f.write(f"- [{i+1}. Segment at {c['timestamp']:.2f}s](#segment-{i+1})\n")
            
            f.write("\n---\n\n")
            
            for i, c in enumerate(concepts):
                f.write(f"### Segment {i+1}\n")
                f.write(f"*Timestamp: {c['timestamp']:.2f}s*\n\n")
                f.write("#### Original Transcript Snippet\n")
                f.write(f"> {c['text']}\n\n")
                f.write("#### Synthesized Knowledge\n")
                f.write(f"{c.get('synthesized_article', 'No synthesis available.')}\n\n")
                f.write("---\n\n")
        
        return filepath

class PPTGenerator:
    """
    Generates PowerPoint slides (PPTX) summarizing the extracted logic.
    """
    def __init__(self, output_dir="data/outputs"):
        from pptx import Presentation
        self.output_dir = output_dir
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

    def generate(self, video_title, youtube_url, richness_score, concepts):
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation()
        
        # Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = f"V2K Sentinel Analysis:\n{video_title}"
        subtitle.text = f"URL: {youtube_url}\nRichness Score: {richness_score:.2f}/1.00\nGenerated: {datetime.now().strftime('%Y-%m-%d')}"

        # Concept Slides
        bullet_slide_layout = prs.slide_layouts[1]
        for i, concept in enumerate(concepts):
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            title_shape.text = f"Segment {i+1} (at {concept['timestamp']:.2f}s)"
            
            tf = body_shape.text_frame
            tf.text = concept.get('synthesized_article', 'No synthesis available.')[:500] + "..." # Limit text for slide

        safe_title = "".join([c for c in video_title if c.isalnum() or c in (' ', '_')]).rstrip()
        filename = f"{safe_title}.pptx"
        filepath = os.path.join(self.output_dir, filename)
        prs.save(filepath)
        return filepath

if __name__ == "__main__":
    test_concepts = [
        {"timestamp": 0.0, "text": "Hello world", "synthesized_article": "Intro knowledge article content..."},
        {"timestamp": 30.5, "text": "Deep learning intro", "synthesized_article": "Neural networks are models..."}
    ]
    gen = PDFGenerator()
    path = gen.generate("Test Video", "http://test.com", 0.95, test_concepts)
    print(f"Generated PDF at: {path}")
